import os
import re
import shutil
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Set
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
import docx
import pandas as pd  # Import pandas to read Excel and CSV files
from pptx import Presentation  # Import Presentation for PPT files

# Legal document extensions and patterns
LEGAL_DOC_EXTENSIONS = {
    '.motion', '.pleading', '.affidavit', '.declaration', 
    '.brief', '.memorandum', '.order', '.judgment', '.opinion',
    '.memo', '.memoofpanda', '.memoofp&a', '.mtv', '.msj', '.mtd'
}

# Common legal document patterns for content-based detection
LEGAL_DOC_PATTERNS = [
    r'IN THE (?:\w+ )*COURT',
    r'CASE NO\.?\s*[:\[].*\d',
    r'PLAINTIFF.*V\.?\s*DEFENDANT',
    r'COMPLAINANT.*V\.?\s*RESPONDENT',
    r'IN RE:',
    r'MOTION TO',
    r'MEMORANDUM OF (LAW|POINTS AND AUTHORITIES)',
    r'UNITED STATES DISTRICT COURT',
    r'SUPREME COURT',
    r'COURT OF APPEALS',
    r'JUDICIAL DISTRICT',
    r'CIRCUIT COURT',
    r'COUNTY COURT',
    r'BEFORE THE HONORABLE',
    r'JUDGE PRESIDING'
]

# Text file extensions
TEXT_EXTENSIONS = {
    '.txt', '.md', '.markdown', '.rst', '.json', '.xml', 
    '.html', '.htm', '.css', '.js', '.py', '.java', '.c', 
    '.cpp', '.h', '.hpp', '.cs', '.php', '.rb', '.go', 
    '.rs', '.swift', '.kt', '.ts', '.jsx', '.tsx', '.log'
}

# Image file extensions
IMAGE_EXTENSIONS = {
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp',
    '.svg', '.heic', '.heif', '.raw', '.cr2', '.nef', '.arw'
}

# Document file extensions
DOCUMENT_EXTENSIONS = {
    '.pdf', '.doc', '.docx', '.odt', '.rtf', '.tex', '.wpd',
    '.pages', '.numbers', '.key', '.odp', '.pps', '.ppsx'
}

# Spreadsheet file extensions
SPREADSHEET_EXTENSIONS = {
    '.xls', '.xlsx', '.xlsm', '.xlsb', '.ods', '.csv', '.tsv'
}

# Presentation file extensions
PRESENTATION_EXTENSIONS = {
    '.ppt', '.pptx', '.odp', '.key', '.pps', '.ppsx'
}

# Archive file extensions
ARCHIVE_EXTENSIONS = {
    '.zip', '.rar', '.7z', '.tar', '.gz', '.bz2', '.xz',
    '.iso', '.dmg', '.pkg', '.deb', '.rpm'
}

# Media file extensions
MEDIA_EXTENSIONS = {
    '.mp3', '.wav', '.ogg', '.flac', '.aac', '.m4a',
    '.mp4', '.avi', '.mov', '.wmv', '.flv', '.mkv',
    '.m4v', '.webm', '.mpg', '.mpeg'
}

def is_legal_document(file_path: str, content: Optional[str] = None) -> bool:
    """Determine if a file is a legal document.
    
    Args:
        file_path: Path to the file
        content: Optional content of the file (if already read)
        
    Returns:
        bool: True if the file appears to be a legal document
    """
    # Check file extension
    ext = os.path.splitext(file_path)[1].lower()
    if ext in LEGAL_DOC_EXTENSIONS:
        return True
        
    # Check content if provided
    if content:
        # Look for common legal document patterns
        content_upper = content.upper()
        for pattern in LEGAL_DOC_PATTERNS:
            if re.search(pattern, content_upper):
                return True
                
    return False

def get_file_category(file_path: str, content: Optional[str] = None) -> str:
    """Determine the category of a file based on its extension and optionally content.
    
    Args:
        file_path: Path to the file
        content: Optional content of the file (if already read)
        
    Returns:
        str: Category of the file ('legal', 'image', 'text', 'document', 'spreadsheet', 
             'presentation', 'archive', 'media', 'other')
    """
    # Check for legal documents first
    if is_legal_document(file_path, content):
        return 'legal'
    
    # Check by extension
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext in IMAGE_EXTENSIONS:
        return 'image'
    elif ext in TEXT_EXTENSIONS:
        return 'text'
    elif ext in DOCUMENT_EXTENSIONS:
        return 'document'
    elif ext in SPREADSHEET_EXTENSIONS:
        return 'spreadsheet'
    elif ext in PRESENTATION_EXTENSIONS:
        return 'presentation'
    elif ext in ARCHIVE_EXTENSIONS:
        return 'archive'
    elif ext in MEDIA_EXTENSIONS:
        return 'media'
    else:
        return 'other'

def separate_files_by_type(file_paths: List[str], check_content: bool = False) -> Dict[str, List[str]]:
    """Separate files into categories based on their extensions and optionally content.
    
    Args:
        file_paths: List of file paths to categorize
        check_content: If True, will read file content to detect legal documents
                      that might not have standard extensions
        
    Returns:
        Dictionary mapping categories to lists of file paths
    """
    categories = {
        'legal': [],
        'images': [],
        'text_files': [],
        'documents': [],
        'spreadsheets': [],
        'presentations': [],
        'archives': [],
        'media': [],
        'others': []
    }
    
    for file_path in file_paths:
        try:
            content = None
            if check_content and os.path.getsize(file_path) < 10 * 1024 * 1024:  # Limit to 10MB
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read(8192)  # Only read first 8KB for detection
                except:
                    pass
            
            category = get_file_category(file_path, content)
            
            # Map the category to the appropriate key in the categories dict
            if category == 'legal':
                categories['legal'].append(file_path)
            elif category == 'image':
                categories['images'].append(file_path)
            elif category == 'text':
                categories['text_files'].append(file_path)
            elif category == 'document':
                categories['documents'].append(file_path)
            elif category == 'spreadsheet':
                categories['spreadsheets'].append(file_path)
            elif category == 'presentation':
                categories['presentations'].append(file_path)
            elif category == 'archive':
                categories['archives'].append(file_path)
            elif category == 'media':
                categories['media'].append(file_path)
            else:
                categories['others'].append(file_path)
                
        except Exception as e:
            print(f"Error processing file {file_path}: {e}")
            categories['others'].append(file_path)
    
    return categories

def read_text_file(file_path):
    """Read text content from a text file."""
    max_chars = 3000  # Limit processing time
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read(max_chars)
        return text
    except Exception as e:
        print(f"Error reading text file {file_path}: {e}")
        return None

def read_docx_file(file_path):
    """Read text content from a .docx or .doc file."""
    try:
        doc = docx.Document(file_path)
        full_text = [para.text for para in doc.paragraphs]
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Error reading DOCX file {file_path}: {e}")
        return None

def read_pdf_file(file_path):
    """Read text content from a PDF file."""
    try:
        doc = fitz.open(file_path)
        # Read only the first few pages to speed up processing
        num_pages_to_read = 3  # Adjust as needed
        full_text = []
        for page_num in range(min(num_pages_to_read, len(doc))):
            page = doc.load_page(page_num)
            full_text.append(page.get_text())
        pdf_content = '\n'.join(full_text)
        return pdf_content
    except Exception as e:
        print(f"Error reading PDF file {file_path}: {e}")
        return None

def read_spreadsheet_file(file_path):
    """Read text content from an Excel or CSV file."""
    try:
        if file_path.lower().endswith('.csv'):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
        text = df.to_string()
        return text
    except Exception as e:
        print(f"Error reading spreadsheet file {file_path}: {e}")
        return None

def read_ppt_file(file_path):
    """Read text content from a PowerPoint file."""
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Error reading PowerPoint file {file_path}: {e}")
        return None

def read_file_data(file_path: str, max_size: int = 10 * 1024 * 1024) -> Optional[str]:
    """Read content from a file based on its extension.
    
    Args:
        file_path: Path to the file to read
        max_size: Maximum file size to read in bytes (default: 10MB)
        
    Returns:
        str: The file content, or None if the file could not be read
    """
    try:
        # Check file size first
        file_size = os.path.getsize(file_path)
        if file_size > max_size:
            print(f"File {file_path} is too large ({file_size/1024/1024:.2f}MB > {max_size/1024/1024}MB)")
            return None
            
        _, ext = os.path.splitext(file_path.lower())
        
        # Try to detect if it's a legal document first
        if is_legal_document(file_path, None):
            # For legal documents, read as text first
            content = read_text_file(file_path)
            if content and any(pattern in content.upper() for pattern in LEGAL_DOC_PATTERNS):
                return content
        
        # Fall back to extension-based reading
        if ext in TEXT_EXTENSIONS:
            return read_text_file(file_path)
        elif ext in {'.docx', '.doc'}:
            return read_docx_file(file_path)
        elif ext == '.pdf':
            return read_pdf_file(file_path)
        elif ext in SPREADSHEET_EXTENSIONS:
            return read_spreadsheet_file(file_path)
        elif ext in PRESENTATION_EXTENSIONS:
            return read_ppt_file(file_path)
        else:
            # Try to read as text file as a fallback
            try:
                return read_text_file(file_path)
            except:
                return None
    except Exception as e:
        print(f"Error reading file {file_path}: {e}")
        return None

def display_directory_tree(path):
    """Display the directory tree in a format similar to the 'tree' command, including the full path."""
    def tree(dir_path, prefix=''):
        contents = sorted([c for c in os.listdir(dir_path) if not c.startswith('.')])
        pointers = ['├── '] * (len(contents) - 1) + ['└── '] if contents else []
        for pointer, name in zip(pointers, contents):
            full_path = os.path.join(dir_path, name)
            print(prefix + pointer + name)
            if os.path.isdir(full_path):
                extension = '│   ' if pointer == '├── ' else '    '
                tree(full_path, prefix + extension)
    if os.path.isdir(path):
        print(os.path.abspath(path))
        tree(path)
    else:
        print(os.path.abspath(path))

def collect_file_paths(base_path):
    """Collect all file paths from the base directory or single file, excluding hidden files."""
    if os.path.isfile(base_path):
        return [base_path]
    else:
        file_paths = []
        for root, _, files in os.walk(base_path):
            for file in files:
                if not file.startswith('.'):  # Exclude hidden files
                    file_paths.append(os.path.join(root, file))
        return file_paths

def separate_files_by_type(file_paths):
    """Separate files into images and text files based on their extensions."""
    image_extensions = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff')
    text_extensions = ('.txt', '.docx', '.doc', '.pdf', '.md', '.xls', '.xlsx', '.ppt', '.pptx', '.csv')
    image_files = [fp for fp in file_paths if os.path.splitext(fp.lower())[1] in image_extensions]
    text_files = [fp for fp in file_paths if os.path.splitext(fp.lower())[1] in text_extensions]

    return image_files, text_files  # Return only two values

# TODO:ebook: '.mobi', '.azw', '.azw3', '.epub',
